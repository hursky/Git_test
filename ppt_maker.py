from collections import abc
from pptx import Presentation
from pptx.util import Cm, Pt, Mm, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from analyze import run_analysis, ALGORITHM_VER
import os
from PySide6.QtWidgets import QProgressBar, QStatusBar

def make_ppt(conf, save_dir: str = None, progressbar:QProgressBar=None, statusbar:QStatusBar=None):
    if save_dir is None:
        save_path = f"./{conf['site']}_구축신호_분석_보고서_{conf['date']}.pptx"
    else:
        save_path = os.path.join(save_dir, f"{conf['site']}_구축신호_분석_보고서_{conf['date']}.pptx")
    # 프레젠테이션
    statusbar.showMessage(f"ppt 초기화")
    prs = Presentation()
    prs.slide_width = Cm(33.87)
    prs.slide_height = Cm(19.05)

    front_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(front_layout)
    txBox = slide.shapes.add_textbox(Cm(6.14), Cm(
        3.39), Cm(21.59), Cm(4.08))  # left, top, width, height
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{conf['site']} \n구축 신호 분석 보고서"
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    txBox = slide.shapes.add_textbox(Cm(6.14), Cm(
        10.8), Cm(21.59), Cm(4.08))  # left, top, width, height
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"구축 엔지니어: {conf['engineer']}\n분석일: {conf['date']}\n분석 알고리즘 버전: {ALGORITHM_VER}"
    p.font.size = Pt(30)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(108, 110, 105)

    list_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(list_layout)
    title = slide.placeholders[0]
    title.text = "대상 모터 목록"
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    m_set = []
    for m in conf["motor_set"]:
        if (m["name"] in [None, ""]) or len(m['data']) == 0:
            continue
        m_set.append(m)

    for m in m_set:
        p = tf.add_paragraph()
        p.text = m["name"]
        p.font.size = Pt(20)
        p.font.bold = True

    # 분석 보고서
    progressbar.setMaximum(len(m_set))
    progressbar.setValue(0)
    for pi, m in enumerate(m_set):
        progressbar.setValue(pi+1)
        statusbar.showMessage(f"{m['name']} 분석중..")
        ret = run_analysis(motor_name=m["name"], aws_files=m["data"])
        # 메인 분석 페이지 추가
        report_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(report_layout)

        # 상단 제목
        txBox = slide.shapes.add_textbox(Cm(2.33), Cm(
            1), Cm(29), Cm(1.5))  # left, top, width, height
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{m['name']} 신호 분석 - 운전신호"
        p.font.size = Pt(20)
        p.font.bold = True

        # 운전신호
        tf = slide.shapes.add_textbox(left=Cm(7.85), top=Cm(
            2.5), width=Cm(2.62), height=Cm(0.85)).text_frame
        p = tf.paragraphs[0]
        p.text = "운전신호 (RST 상)"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        pic = slide.shapes.add_picture(
            ret["driving_rst_pic"], left=Cm(2.33), top=Cm(3.36), width=Cm(13.67))

        # FFT 분석
        tf = slide.shapes.add_textbox(left=Cm(23.54), top=Cm(
            2.5), width=Cm(2.62), height=Cm(0.85)).text_frame
        p = tf.paragraphs[0]
        p.text = "FFT 분석"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        pic = slide.shapes.add_picture(
            ret["driving_fft_pic"], left=Cm(17.87), top=Cm(3.36), width=Cm(13.67))

        # 텍스트
        tf = slide.shapes.add_textbox(left=Cm(2.33), top=Cm(
            10.65), width=Cm(21.22), height=Cm(7)).text_frame
        p = tf.paragraphs[0]
        p.text = ret["driving_text"]
        p.font.size = Pt(10)

        # 순간 최대 변동폭 (RMS-A)
        tf = slide.shapes.add_textbox(left=Cm(26.55), top=Cm(
            10.65), width=Cm(3.12), height=Cm(0.85)).text_frame
        p = tf.paragraphs[0]
        p.text = "순간 최대 변동폭 (RMS-A)"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        pic = slide.shapes.add_picture(
            ret["driving_p2p_max_pic"], left=Cm(24.67), top=Cm(11.5), width=Cm(6.87))

        # 추가 분석 페이지 추가
        report_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(report_layout)

        # 상단 제목
        txBox = slide.shapes.add_textbox(Cm(2.33), Cm(
            1), Cm(29), Cm(1.5))  # left, top, width, height
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{m['name']} 신호 분석 - 기타 신호"
        p.font.size = Pt(20)
        p.font.bold = True

        # 기동 순간 신호
        tf = slide.shapes.add_textbox(left=Cm(7.85), top=Cm(
            2.5), width=Cm(2.62), height=Cm(0.85)).text_frame
        p = tf.paragraphs[0]
        p.text = "기동 순간 신호"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        pic = slide.shapes.add_picture(
            ret["on_start_pic"], left=Cm(2.33), top=Cm(3.36), width=Cm(13.67))

        # 정지 신호
        tf = slide.shapes.add_textbox(left=Cm(23.54), top=Cm(
            2.5), width=Cm(2.62), height=Cm(0.85)).text_frame
        p = tf.paragraphs[0]
        p.text = "정지 신호"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        pic = slide.shapes.add_picture(
            ret["stop_pic"], left=Cm(17.87), top=Cm(3.36), width=Cm(13.67))

        # 정지 순간 신호
        tf = slide.shapes.add_textbox(left=Cm(7.85), top=Cm(
            10.54), width=Cm(2.62), height=Cm(0.85)).text_frame
        p = tf.paragraphs[0]
        p.text = "정지 순간 신호"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        pic = slide.shapes.add_picture(
            ret["on_stop_pic"], left=Cm(2.33), top=Cm(11.39), width=Cm(13.67))

    prs.save(save_path)
    statusbar.showMessage("결과 리포트 생성 완료!")
    return save_path
